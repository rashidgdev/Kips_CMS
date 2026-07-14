"""
Generates the client presentation deck for the KIPS College Kasur Campus
Management System. One-off script, not part of the Django app - run with:
    python scripts/generate_presentation.py
"""
import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
LOGO_ICON = os.path.join(BASE_DIR, 'static', 'images', 'kips-logo-icon.png')
LOGO_FULL = os.path.join(BASE_DIR, 'static', 'images', 'kips-logo-full.png')
OUT_PATH = os.path.join(BASE_DIR, 'KIPS_CMS_Client_Presentation.pptx')

# ---- Palette (matches the app's own design system) -------------------------
PRIMARY = RGBColor(0x1E, 0x3A, 0x8A)       # blue-900
PRIMARY_MID = RGBColor(0x1D, 0x4E, 0xD8)   # blue-700
INDIGO = RGBColor(0x43, 0x38, 0xCA)        # indigo-700
BG_DARK_FROM = RGBColor(0x17, 0x2B, 0x6B)
BG_DARK_TO = RGBColor(0x31, 0x28, 0x99)
BG_LIGHT = RGBColor(0xF8, 0xFA, 0xFC)
CARD_BG = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BORDER = RGBColor(0xE5, 0xE7, 0xEB)
TEXT_DARK = RGBColor(0x11, 0x18, 0x27)
TEXT_MUTED = RGBColor(0x6B, 0x72, 0x80)
TEXT_ON_DARK = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_ON_DARK_MUTED = RGBColor(0xC7, 0xD2, 0xFE)
GREEN = RGBColor(0x05, 0x96, 0x69)
GREEN_BG = RGBColor(0xD1, 0xFA, 0xE5)
AMBER = RGBColor(0xB4, 0x53, 0x09)
AMBER_BG = RGBColor(0xFE, 0xF3, 0xC7)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

FONT = 'Segoe UI'

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


# ---- low-level helpers ------------------------------------------------------

def add_slide():
    return prs.slides.add_slide(BLANK)


def set_background(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def gradient_rect(slide, left, top, width, height, color_from, color_to, angle=90):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    shape.shadow.inherit = False
    fill = shape.fill
    fill.gradient()
    stops = fill.gradient_stops
    stops[0].color.rgb = color_from
    stops[0].position = 0.0
    stops[1].color.rgb = color_to
    stops[1].position = 1.0
    try:
        fill.gradient_angle = angle
    except Exception:
        pass
    return shape


def rect(slide, left, top, width, height, fill_color=None, line_color=None, line_w=0.75, shadow=False, radius=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius is not None else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.shadow.inherit = False
    if radius is not None:
        try:
            shape.adjustments[0] = radius
        except Exception:
            pass
    if fill_color is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color is not None:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_w)
    else:
        shape.line.fill.background()
    return shape


def textbox(slide, left, top, width, height, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return box, tf


def set_run(run, text, size=18, color=TEXT_DARK, bold=False, italic=False, font=FONT):
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font


def add_para(tf, text, size=18, color=TEXT_DARK, bold=False, italic=False, align=PP_ALIGN.LEFT,
             space_after=6, space_before=0, bullet=False, level=0, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.level = level
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    run = p.add_run()
    set_run(run, text, size=size, color=color, bold=bold, italic=italic)
    if bullet:
        _set_bullet(p, color)
    else:
        _no_bullet(p)
    return p


def _no_bullet(p):
    pPr = p._pPr
    if pPr is None:
        pPr = p._p.get_or_add_pPr()
    buNone = pPr.makeelement(qn('a:buNone'), {})
    pPr.append(buNone)


def _set_bullet(p, color, char='•'):
    pPr = p._p.get_or_add_pPr()
    buFont = pPr.makeelement(qn('a:buFont'), {'typeface': 'Arial'})
    buChar = pPr.makeelement(qn('a:buChar'), {'char': char})
    pPr.append(buFont)
    pPr.append(buChar)


def footer(slide, page_num, label='KIPS College Kasur Campus — Campus Management System'):
    rect(slide, 0, Inches(7.13), SLIDE_W, Inches(0.37), fill_color=BG_LIGHT)
    box, tf = textbox(slide, Inches(0.55), Inches(7.15), Inches(9), Inches(0.32), anchor=MSO_ANCHOR.MIDDLE)
    add_para(tf, label, size=9, color=TEXT_MUTED, first=True)
    box2, tf2 = textbox(slide, Inches(12.4), Inches(7.15), Inches(0.6), Inches(0.32), anchor=MSO_ANCHOR.MIDDLE)
    add_para(tf2, str(page_num), size=9, color=TEXT_MUTED, align=PP_ALIGN.RIGHT, first=True)
    if os.path.exists(LOGO_ICON):
        slide.shapes.add_picture(LOGO_ICON, Inches(12.75), Inches(7.13), height=Inches(0.35))


def slide_header(slide, kicker, title, dark=False):
    kicker_color = INDIGO if not dark else TEXT_ON_DARK_MUTED
    title_color = TEXT_DARK if not dark else TEXT_ON_DARK
    box, tf = textbox(slide, Inches(0.55), Inches(0.35), Inches(11.5), Inches(0.4))
    add_para(tf, kicker.upper(), size=13, color=kicker_color, bold=True, first=True)
    box2, tf2 = textbox(slide, Inches(0.55), Inches(0.68), Inches(12), Inches(0.7))
    add_para(tf2, title, size=30, color=title_color, bold=True, first=True)
    rect(slide, Inches(0.57), Inches(1.42), Inches(0.6), Pt(4), fill_color=INDIGO)


# ---- slide builders ----------------------------------------------------------

def title_slide():
    s = add_slide()
    gradient_rect(s, 0, 0, SLIDE_W, SLIDE_H, BG_DARK_FROM, BG_DARK_TO, angle=45)
    # glow accents
    c1 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(-2), Inches(6), Inches(6))
    c1.fill.solid(); c1.fill.fore_color.rgb = RGBColor(0x3B, 0x82, 0xF6); c1.fill.transparency = 0
    c1.line.fill.background(); c1.shadow.inherit = False
    _set_transparency(c1, 80)
    c2 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.5), Inches(4.5), Inches(6), Inches(6))
    c2.fill.solid(); c2.fill.fore_color.rgb = INDIGO
    c2.line.fill.background(); c2.shadow.inherit = False
    _set_transparency(c2, 80)

    if os.path.exists(LOGO_ICON):
        pic = s.shapes.add_picture(LOGO_ICON, Inches(5.92), Inches(0.75), height=Inches(1.3))

    box, tf = textbox(s, Inches(1), Inches(2.3), Inches(11.33), Inches(0.5), anchor=MSO_ANCHOR.MIDDLE)
    add_para(tf, 'KIPS COLLEGE KASUR CAMPUS', size=16, color=TEXT_ON_DARK_MUTED, bold=True,
             align=PP_ALIGN.CENTER, first=True)

    box2, tf2 = textbox(s, Inches(0.8), Inches(2.75), Inches(11.73), Inches(1.6), anchor=MSO_ANCHOR.MIDDLE)
    add_para(tf2, 'Campus Management System', size=44, color=TEXT_ON_DARK, bold=True,
             align=PP_ALIGN.CENTER, first=True)

    box3, tf3 = textbox(s, Inches(1.5), Inches(4.15), Inches(10.33), Inches(0.6), anchor=MSO_ANCHOR.MIDDLE)
    add_para(tf3, 'A single web portal replacing manual registers and spreadsheets —'
                  ' academics, attendance, faculty payroll, fees, and reporting in one place.',
             size=16, color=TEXT_ON_DARK_MUTED, align=PP_ALIGN.CENTER, first=True)

    box4, tf4 = textbox(s, Inches(1), Inches(6.55), Inches(11.33), Inches(0.5), anchor=MSO_ANCHOR.MIDDLE)
    add_para(tf4, 'Client Presentation  •  Project Delivery Overview', size=13,
             color=TEXT_ON_DARK_MUTED, align=PP_ALIGN.CENTER, first=True)


def section_divider(page, number, title, subtitle):
    s = add_slide()
    gradient_rect(s, 0, 0, SLIDE_W, SLIDE_H, BG_DARK_FROM, BG_DARK_TO, angle=45)
    box, tf = textbox(s, Inches(1), Inches(2.7), Inches(2), Inches(1))
    add_para(tf, number, size=72, color=RGBColor(0x93, 0xA5, 0xF5), bold=True, first=True)
    box2, tf2 = textbox(s, Inches(1), Inches(3.7), Inches(11.33), Inches(1))
    add_para(tf2, title, size=38, color=TEXT_ON_DARK, bold=True, first=True)
    box3, tf3 = textbox(s, Inches(1), Inches(4.55), Inches(10.5), Inches(0.8))
    add_para(tf3, subtitle, size=16, color=TEXT_ON_DARK_MUTED, first=True)
    footer(s, page)
    return s


def _set_transparency(shape, pct):
    """pct: 0 opaque .. 100 fully transparent"""
    sp = shape.fill.fore_color._xFill
    srgbClr = sp.find(qn('a:srgbClr'))
    if srgbClr is None:
        return
    alpha = srgbClr.makeelement(qn('a:alpha'), {'val': str(int((100 - pct) * 1000))})
    srgbClr.append(alpha)


def bullet_slide(page, kicker, title, bullets, note=None):
    s = add_slide()
    set_background(s, RGBColor(0xFF, 0xFF, 0xFF))
    slide_header(s, kicker, title)
    box, tf = textbox(s, Inches(0.7), Inches(1.85), Inches(11.7), Inches(4.7))
    first = True
    for b in bullets:
        if isinstance(b, tuple):
            head, sub = b
            add_para(tf, head, size=18, color=TEXT_DARK, bold=True, bullet=True,
                     space_after=2, first=first)
            add_para(tf, sub, size=14, color=TEXT_MUTED, level=1, space_after=14, first=False)
        else:
            add_para(tf, b, size=18, color=TEXT_DARK, bullet=True, space_after=12, first=first)
        first = False
    if note:
        rect(s, Inches(0.7), Inches(6.35), Inches(11.7), Inches(0.65), fill_color=RGBColor(0xEE, 0xF2, 0xFF), radius=0.25)
        nb, ntf = textbox(s, Inches(1.0), Inches(6.35), Inches(11.1), Inches(0.65), anchor=MSO_ANCHOR.MIDDLE)
        add_para(ntf, note, size=13, color=PRIMARY, italic=True, first=True)
    footer(s, page)
    return s


def two_col_slide(page, kicker, title, left_title, left_items, right_title, right_items,
                   left_color=PRIMARY_MID, right_color=AMBER):
    s = add_slide()
    set_background(s, RGBColor(0xFF, 0xFF, 0xFF))
    slide_header(s, kicker, title)

    col_w = Inches(5.7)
    for i, (heading, items, accent) in enumerate([
        (left_title, left_items, left_color), (right_title, right_items, right_color)
    ]):
        x = Inches(0.7) + i * (col_w + Inches(0.3))
        rect(s, x, Inches(1.9), col_w, Inches(4.9), fill_color=BG_LIGHT, line_color=CARD_BORDER, radius=0.05)
        rect(s, x, Inches(1.9), Inches(0.09), Inches(4.9), fill_color=accent)
        hb, htf = textbox(s, x + Inches(0.35), Inches(2.15), col_w - Inches(0.6), Inches(0.5))
        add_para(htf, heading, size=18, color=TEXT_DARK, bold=True, first=True)
        ib, itf = textbox(s, x + Inches(0.35), Inches(2.75), col_w - Inches(0.6), Inches(3.9))
        first = True
        for it in items:
            add_para(itf, it, size=14, color=TEXT_DARK, bullet=True, space_after=10, first=first)
            first = False
    footer(s, page)
    return s


def roadmap_slide(page, kicker, title, items):
    """items: list of (label, sublabel) all shown complete."""
    s = add_slide()
    set_background(s, RGBColor(0xFF, 0xFF, 0xFF))
    slide_header(s, kicker, title)

    n = len(items)
    top = Inches(3.05)
    left_margin = Inches(0.7)
    right_margin = Inches(0.7)
    usable = SLIDE_W - left_margin - right_margin
    step = usable / n
    line_y = top + Inches(0.32)
    rect(s, left_margin + step / 2, line_y, usable - step, Pt(3), fill_color=CARD_BORDER)

    for i, (label, sub) in enumerate(items):
        cx = left_margin + step * i + step / 2
        d = Inches(0.64)
        node = s.shapes.add_shape(MSO_SHAPE.OVAL, cx - d / 2, top, d, d)
        node.fill.solid(); node.fill.fore_color.rgb = GREEN
        node.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF); node.line.width = Pt(2)
        node.shadow.inherit = False
        ntf = node.text_frame
        ntf.word_wrap = False
        ntf.margin_left = 0; ntf.margin_right = 0; ntf.margin_top = 0; ntf.margin_bottom = 0
        p = ntf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        set_run(run, '✓', size=20, color=RGBColor(0xFF, 0xFF, 0xFF), bold=True)

        lbl_box, lbl_tf = textbox(s, cx - step / 2 + Inches(0.05), top + Inches(0.85),
                                   step - Inches(0.1), Inches(0.85))
        add_para(lbl_tf, label, size=12.5, color=TEXT_DARK, bold=True, align=PP_ALIGN.CENTER, first=True,
                 space_after=2)
        sub_box, sub_tf = textbox(s, cx - step / 2 + Inches(0.05), top + Inches(1.55),
                                   step - Inches(0.1), Inches(1.6))
        add_para(sub_tf, sub, size=10.5, color=TEXT_MUTED, align=PP_ALIGN.CENTER, first=True)

    tagbox, tagtf = textbox(s, Inches(0.7), Inches(6.25), Inches(11.7), Inches(0.5), anchor=MSO_ANCHOR.MIDDLE)
    add_para(tagtf, 'All phases delivered and verified end-to-end — the system is complete and in testing.',
             size=13, color=GREEN, bold=True, align=PP_ALIGN.CENTER, first=True)
    footer(s, page)
    return s


def roles_slide(page):
    s = add_slide()
    set_background(s, RGBColor(0xFF, 0xFF, 0xFF))
    slide_header(s, 'Access Control', 'Six Roles, Six Tailored Experiences')

    roles = [
        ('Student', 'Attendance %, grades & GPA, timetable, fee status — all self-service.'),
        ('Teacher', 'Mark attendance, log lectures, enter marks, view teaching schedule.'),
        ('Head of Department', 'Everything a teacher has, plus department-level oversight.'),
        ('Campus Coordinator', 'Programs, courses, faculty assignments, enrollments, timetable.'),
        ('Accountant', 'Fee structures, student billing, payment collection, financial reports.'),
        ('College Administrator', 'Full control: every module, every report, user management.'),
    ]
    cols, rows = 3, 2
    card_w = Inches(3.75)
    card_h = Inches(2.05)
    gap = Inches(0.2)
    start_x = Inches(0.7)
    start_y = Inches(1.95)
    for idx, (name, desc) in enumerate(roles):
        r, c = divmod(idx, cols)
        x = start_x + c * (card_w + gap)
        y = start_y + r * (card_h + gap)
        rect(s, x, y, card_w, card_h, fill_color=BG_LIGHT, line_color=CARD_BORDER, radius=0.08)
        rect(s, x, y, card_w, Inches(0.08), fill_color=PRIMARY_MID)
        nb, ntf = textbox(s, x + Inches(0.25), y + Inches(0.22), card_w - Inches(0.5), Inches(0.45))
        add_para(ntf, name, size=15.5, color=TEXT_DARK, bold=True, first=True)
        db, dtf = textbox(s, x + Inches(0.25), y + Inches(0.72), card_w - Inches(0.5), Inches(1.2))
        add_para(dtf, desc, size=12, color=TEXT_MUTED, first=True)
    footer(s, page)
    return s


def module_slide(page, number, title, tagline, bullets, benefit):
    s = add_slide()
    set_background(s, RGBColor(0xFF, 0xFF, 0xFF))

    # left accent panel
    rect(s, 0, 0, Inches(3.6), SLIDE_H, fill_color=PRIMARY)
    nb, ntf = textbox(s, Inches(0.5), Inches(0.55), Inches(2.6), Inches(1), anchor=MSO_ANCHOR.TOP)
    add_para(ntf, f'MODULE {number}', size=13, color=TEXT_ON_DARK_MUTED, bold=True, first=True)
    tb, ttf = textbox(s, Inches(0.5), Inches(1.0), Inches(2.9), Inches(2.2))
    add_para(ttf, title, size=26, color=TEXT_ON_DARK, bold=True, first=True)
    tgb, tgtf = textbox(s, Inches(0.5), Inches(2.55), Inches(2.9), Inches(1.6))
    add_para(tgtf, tagline, size=13, color=TEXT_ON_DARK_MUTED, first=True)

    rect(s, Inches(0.5), Inches(5.55), Inches(2.6), Pt(2), fill_color=RGBColor(0x5B, 0x6E, 0xC9))
    bb, btf = textbox(s, Inches(0.5), Inches(5.8), Inches(2.9), Inches(1.2))
    add_para(btf, 'CLIENT BENEFIT', size=11, color=TEXT_ON_DARK_MUTED, bold=True, first=True, space_after=4)
    add_para(btf, benefit, size=12.5, color=TEXT_ON_DARK, first=False)

    box, tf = textbox(s, Inches(4.0), Inches(0.85), Inches(8.7), Inches(6.0))
    first = True
    for b in bullets:
        if isinstance(b, tuple):
            head, sub = b
            add_para(tf, head, size=16.5, color=TEXT_DARK, bold=True, bullet=True, space_after=2, first=first)
            add_para(tf, sub, size=12.5, color=TEXT_MUTED, level=1, space_after=14, first=False)
        else:
            add_para(tf, b, size=16.5, color=TEXT_DARK, bullet=True, space_after=13, first=first)
        first = False
    footer(s, page)
    return s


def benefits_slide(page):
    s = add_slide()
    set_background(s, RGBColor(0xFF, 0xFF, 0xFF))
    slide_header(s, 'Value Delivered', 'Why This Matters for the Campus')

    items = [
        ('Zero Manual Registers', 'Attendance, day books, and result sheets are no longer paper- or '
                                   'spreadsheet-based — one system, always current.'),
        ('Real-Time Visibility', 'Students see attendance/grades/fees the moment they are entered; '
                                  'management sees campus-wide numbers instantly, not at month-end.'),
        ('Payroll Tied to Real Data', "Faculty pay is calculated from verified lecture counts, not "
                                       "estimated — removing a recurring source of dispute."),
        ('No Double-Booking', 'Timetable scheduling blocks teacher and room conflicts automatically '
                               'before they happen.'),
        ('Full Financial Trail', 'Every fee item and payment is tracked per student, per semester, '
                                  'with instant outstanding-balance reporting.'),
        ('One Login, One Dashboard', 'Every role — student to administrator — gets exactly the '
                                      'tools relevant to them, nothing more, nothing less.'),
    ]
    cols = 2
    card_w = Inches(5.7)
    card_h = Inches(1.42)
    gap_x = Inches(0.3)
    gap_y = Inches(0.18)
    start_x = Inches(0.7)
    start_y = Inches(1.85)
    for idx, (h, d) in enumerate(items):
        r, c = divmod(idx, cols)
        x = start_x + c * (card_w + gap_x)
        y = start_y + r * (card_h + gap_y)
        rect(s, x, y, card_w, card_h, fill_color=BG_LIGHT, line_color=CARD_BORDER, radius=0.12)
        chip = rect(s, x + Inches(0.2), y + Inches(0.2), Inches(0.42), Inches(0.42), fill_color=GREEN_BG, radius=0.5)
        cftf = chip.text_frame
        cftf.margin_left = 0; cftf.margin_right = 0; cftf.margin_top = 0; cftf.margin_bottom = 0
        cp = cftf.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        crun = cp.add_run()
        set_run(crun, '✓', size=16, color=GREEN, bold=True)
        hb, htf = textbox(s, x + Inches(0.78), y + Inches(0.16), card_w - Inches(1.0), Inches(0.35))
        add_para(htf, h, size=14.5, color=TEXT_DARK, bold=True, first=True)
        db, dtf = textbox(s, x + Inches(0.78), y + Inches(0.53), card_w - Inches(1.0), Inches(0.85))
        add_para(dtf, d, size=11, color=TEXT_MUTED, first=True)
    footer(s, page)
    return s


def stakeholder_table_slide(page):
    s = add_slide()
    set_background(s, RGBColor(0xFF, 0xFF, 0xFF))
    slide_header(s, 'Value Delivered', 'What Changes for Each Stakeholder')

    rows_data = [
        ('Students', 'Chase teachers for marks; guess attendance status; visit accounts office for fee status',
         'Real-time grades, GPA/CGPA, attendance %, and fee balance from any browser'),
        ('Teachers', 'Paper attendance registers; manual day-book entries; disputed lecture counts',
         'One click to mark attendance and log lectures; pay computed from verified records'),
        ('Coordinators', 'Excel sheets for course/faculty assignment; manual timetable clash-checking',
         'Guided screens for every setup task; conflicts blocked automatically'),
        ('Accountants', 'Manual fee ledgers per student; no consolidated outstanding-balance view',
         'Fee packages generated per student; live due/paid/outstanding dashboard'),
        ('Administration', 'No single source of truth; reports assembled manually for meetings',
         'Campus-wide dashboards and one-click Excel/PDF reports, always current'),
    ]

    headers = ['Stakeholder', 'Before (Manual Process)', 'After (This System)']
    col_w = [Inches(2.0), Inches(4.85), Inches(4.85)]
    start_x = Inches(0.7)
    start_y = Inches(1.85)
    row_h = Inches(0.95)
    header_h = Inches(0.5)

    x = start_x
    for w, htext in zip(col_w, headers):
        rect(s, x, start_y, w, header_h, fill_color=PRIMARY)
        hb, htf = textbox(s, x + Inches(0.15), start_y, w - Inches(0.3), header_h, anchor=MSO_ANCHOR.MIDDLE)
        add_para(htf, htext, size=13, color=TEXT_ON_DARK, bold=True, first=True)
        x += w

    y = start_y + header_h
    for i, (stakeholder, before, after) in enumerate(rows_data):
        row_bg = RGBColor(0xFF, 0xFF, 0xFF) if i % 2 == 0 else BG_LIGHT
        x = start_x
        for w in col_w:
            rect(s, x, y, w, row_h, fill_color=row_bg, line_color=CARD_BORDER, line_w=0.5)
            x += w
        b1, t1 = textbox(s, start_x + Inches(0.15), y, col_w[0] - Inches(0.3), row_h, anchor=MSO_ANCHOR.MIDDLE)
        add_para(t1, stakeholder, size=12.5, color=TEXT_DARK, bold=True, first=True)
        b2, t2 = textbox(s, start_x + col_w[0] + Inches(0.15), y, col_w[1] - Inches(0.3), row_h, anchor=MSO_ANCHOR.MIDDLE)
        add_para(t2, before, size=11, color=AMBER, first=True)
        b3, t3 = textbox(s, start_x + col_w[0] + col_w[1] + Inches(0.15), y, col_w[2] - Inches(0.3), row_h, anchor=MSO_ANCHOR.MIDDLE)
        add_para(t3, after, size=11, color=GREEN, first=True)
        y += row_h

    footer(s, page)
    return s


def closing_slide(page):
    s = add_slide()
    gradient_rect(s, 0, 0, SLIDE_W, SLIDE_H, BG_DARK_FROM, BG_DARK_TO, angle=45)
    if os.path.exists(LOGO_ICON):
        s.shapes.add_picture(LOGO_ICON, Inches(6.17), Inches(0.8), height=Inches(1.1))
    box, tf = textbox(s, Inches(1), Inches(2.35), Inches(11.33), Inches(1), anchor=MSO_ANCHOR.MIDDLE)
    add_para(tf, 'Thank You', size=44, color=TEXT_ON_DARK, bold=True, align=PP_ALIGN.CENTER, first=True)
    box2, tf2 = textbox(s, Inches(1.5), Inches(3.3), Inches(10.33), Inches(0.6), anchor=MSO_ANCHOR.MIDDLE)
    add_para(tf2, 'Questions, feedback, and requested changes are all welcome —'
                  ' this system was built to be shaped around how the campus actually works.',
             size=16, color=TEXT_ON_DARK_MUTED, align=PP_ALIGN.CENTER, first=True)

    rect(s, Inches(4.67), Inches(4.5), Inches(4), Pt(2), fill_color=RGBColor(0x5B, 0x6E, 0xC9))

    box3, tf3 = textbox(s, Inches(1), Inches(4.85), Inches(11.33), Inches(1.5), anchor=MSO_ANCHOR.MIDDLE)
    add_para(tf3, 'KIPS College Kasur Campus', size=16, color=TEXT_ON_DARK, bold=True,
             align=PP_ALIGN.CENTER, first=True, space_after=4)
    add_para(tf3, 'Campus Management System — Project Delivery Team', size=13,
             color=TEXT_ON_DARK_MUTED, align=PP_ALIGN.CENTER, first=False)
    return s


# =============================================================================
# BUILD THE DECK
# =============================================================================

page = 1
title_slide()

page += 1
bullet_slide(
    page, 'Agenda', "What We'll Cover Today",
    [
        'The challenge — what the campus was dealing with before this project',
        'What was asked for — the original requirements',
        'What was delivered — the complete module roadmap',
        'Who uses it — six role-based portals',
        'A closer look at each module',
        'The value this creates — for students, staff, and management',
        "What's next"
    ],
)

page += 1
bullet_slide(
    page, 'The Starting Point', 'A Campus Running on Registers and Spreadsheets',
    [
        ('Attendance & lecture records kept on paper', 'Slow to compile, easy to lose, impossible to search or report on instantly.'),
        ('Faculty pay calculated by hand from those same registers', 'A recurring source of disputes and delay every payroll cycle.'),
        ('Marks, results, and GPA tracked in scattered spreadsheets', 'No single, trustworthy view of a student’s standing.'),
        ('Timetables built manually', 'Teacher and room double-bookings caught only after they happen.'),
        ('Fee records kept per-student on paper ledgers', 'No campus-wide view of who owes what, or when it’s due.'),
        ('Multiple programs and semesters running at once', 'All of the above multiplied across every program, every semester, every term.'),
    ],
    note='This is the exact problem the Campus Management System was commissioned to solve.',
)

page += 1
two_col_slide(
    page, 'Requirements', 'What Was Asked For',
    'Core Modules Requested', [
        'Course & semester management, with faculty assignment',
        'Attendance tracking with shortage warnings',
        'Electronic faculty day book tied to payroll',
        'Student academic portal (marks, GPA/CGPA)',
        'Teacher portal for grading & attendance',
        'Conflict-free timetable scheduling',
        'Fee & financial management per student',
        'Role-specific dashboards & exportable reports',
    ],
    'Non-Negotiables', [
        'Six distinct roles, each seeing only what’s relevant to them',
        'Clean, maintainable, mobile-responsive interface',
        'Input validation and error handling throughout',
        'Sample data for immediate testing',
        'A working local setup with clear run instructions',
        'Built incrementally, reviewed at every stage',
    ],
)

page += 1
roadmap_slide(
    page, 'Delivery', 'Complete Project Roadmap',
    [
        ('Foundation & Access', 'Accounts, roles, and permissions for all 6 user types'),
        ('Attendance', 'Lecture-wise attendance with live % and shortage alerts'),
        ('Faculty Day Book', 'Lecture logging tied directly to payroll'),
        ('Assessments', 'Marks, results, GPA/CGPA, live for students'),
        ('Timetable', 'Conflict-checked scheduling, no double-booking'),
        ('Finance', 'Per-student fee packages, payments, balances'),
        ('Dashboards & Reports', 'Role dashboards, Excel/PDF exports'),
        ('Admin Portal', 'Every setup task self-service, no technical access needed'),
    ],
)

page += 1
roles_slide(page)

page += 1
module_slide(
    page, 1, 'Academic Setup & Administration',
    'The foundation every other module depends on — and now fully self-service.',
    [
        ('Departments, Programs, Semesters, Courses', 'Defined once, reused every term — no re-entry.'),
        ('Faculty Assignments', 'Assign a course + section to a teacher for a semester in one screen.'),
        ('Student & Teacher Onboarding', 'Add a person and their login is created together — one step, not two systems.'),
        ('Enrollments', 'Enroll a student into a course offering; it appears on their dashboard immediately.'),
        ('No Django-admin required', 'Every one of the above is a guided, branded screen — not a raw developer tool.'),
    ],
    'Coordinators and administrators run the entire academic setup themselves, with no IT dependency for routine changes.',
)

page += 1
module_slide(
    page, 2, 'Attendance Management',
    'Subject-wise attendance, marked once, visible everywhere instantly.',
    [
        ('Lecture-wise marking', 'Teacher creates a session and marks the whole class in one screen.'),
        ('Automatic calculations', 'Attendance %, lectures delivered/attended/absent — computed live, never stale.'),
        ('Shortage warnings', 'Students flagged automatically once they fall below the required threshold.'),
        ('Student self-service', 'Every student sees their own daily attendance status, per course.'),
    ],
    'Removes the single biggest paper-based process on campus, with instant visibility for students and staff alike.',
)

page += 1
module_slide(
    page, 3, 'Faculty Day Book & Workload',
    'Every lecture delivered becomes a payroll-ready record automatically.',
    [
        ('Auto-logged from Attendance', 'No duplicate data entry — marking attendance creates the day book entry.'),
        ('Coordinator verification', 'A simple queue to confirm lectures actually took place before they count for pay.'),
        ('Workload & payroll report', 'Delivered vs. verified lecture counts, computed pay, per teacher, per month.'),
        ('Exportable snapshots', 'One-click Excel/PDF for handing straight to payroll.'),
    ],
    'Faculty pay is calculated from verified data, not estimates — removing a recurring source of dispute.',
)

page += 1
module_slide(
    page, 4, 'Assessments & Academic Portal',
    'Quizzes to GPA, calculated the moment a mark is entered.',
    [
        ('Flexible assessment types', 'Quizzes, assignments, presentations, midterms — each with its own weight.'),
        ('Teacher marks entry', 'One screen per assessment, per course, for the whole class.'),
        ('Live grading', 'Percentage, letter grade, semester GPA, and CGPA recalculate instantly on save.'),
        ('Student transparency', 'Every student sees their own obtained/total marks and grade per course, live.'),
    ],
    'Students always know exactly where they stand — no waiting for a result sheet to be compiled by hand.',
)

page += 1
module_slide(
    page, 5, 'Intelligent Timetable',
    'Class scheduling that blocks conflicts before they happen.',
    [
        ('Visual weekly grid', 'Coordinators place a course into a day/room/time slot with a live conflict check.'),
        ('No teacher double-booking', 'The system blocks a teacher being scheduled in two places at once.'),
        ('No room double-booking', 'Same protection applied to every classroom and lab.'),
        ('Everyone sees their own schedule', 'Teachers and students get a personal read-only weekly view.'),
    ],
    'Scheduling conflicts that used to surface on the first day of class are caught before the timetable is even published.',
)

page += 1
module_slide(
    page, 6, 'Fee & Financial Management',
    'Every student’s fee history, tracked from admission to graduation.',
    [
        ('Per-program fee packages', 'Tuition, one-time registration, recurring exam fees — defined once per program.'),
        ('Automatic fee generation', 'One click generates a student’s fee items for their current semester.'),
        ('Payment recording', 'Partial payments supported; outstanding balance always accurate, always current.'),
        ('Status at a glance', 'Paid, partial, overdue, or unpaid — visible instantly to both accountant and student.'),
    ],
    'A complete, always-current financial picture per student — no more reconciling paper ledgers at term-end.',
)

page += 1
module_slide(
    page, 7, 'Dashboards & Reporting',
    'Every number a decision-maker needs, in one place.',
    [
        ('Role-specific dashboards', 'Each of the six roles sees the numbers relevant to them the moment they log in.'),
        ('Attendance & academic reports', 'Per-course gradesheets and attendance sheets, on demand.'),
        ('Semester merit list', 'Every student ranked by GPA, campus-wide.'),
        ('Excel & PDF export, everywhere', 'Every report — workload, financial, academic — exports in one click.'),
    ],
    'Reports that used to take a staff member a day to compile by hand are now available instantly, in the format needed.',
)

page += 1
benefits_slide(page)

page += 1
stakeholder_table_slide(page)

page += 1
bullet_slide(
    page, "What's Next", 'Roadmap Beyond This Delivery',
    [
        ('User acceptance testing', 'Client team tests the live system with real campus data before go-live.'),
        ('Deployment to production hosting', 'Moving from local demonstration to a live, always-on server.'),
        ('Staff onboarding', 'Short walkthroughs for coordinators, accountants, and teachers on their specific screens.'),
        ('Optional enhancements', 'SMS/email notifications, a self-service "change password" flow, and any '
                                   'campus-specific requests that come out of testing.'),
    ],
    note='The system is feature-complete against the original brief and ready for client review today.',
)

closing_slide(page + 1)

prs.save(OUT_PATH)
print(f'Saved: {OUT_PATH}')
print(f'Total slides: {len(prs.slides.__iter__.__self__._sldIdLst)}')
